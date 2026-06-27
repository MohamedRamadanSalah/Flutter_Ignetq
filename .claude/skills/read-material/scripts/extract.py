"""Extract lecture material (PDF / PPTX / DOCX) into a readable bundle.

Usage:
    python extract.py <input-file> [--outdir DIR] [--dpi 150] [--extract-figures]

Output (under <outdir>/<stem>/):
    manifest.json   per-page/slide stats and flags (LOW_TEXT pages MUST be read visually)
    text.md         all extracted text with page/slide markers
    pages/p001.png  raster image of every page/slide (for visual reading)
    figures/        extracted embedded images (PNG/JPG) from PDF/PPTX/DOCX
    figures.json    index of extracted figures with page/slide references

Exit code 0 on success. Prints the manifest path on the last line.
"""
import argparse
import json
import os
import sys
from pathlib import Path


def _save_figure(data: bytes, name: str, figures_dir: Path, fmt: str = "png") -> str:
    """Save a figure to disk, deduplicating by SHA256 content hash."""
    import hashlib
    h = hashlib.sha256(data).hexdigest()[:12]
    fname = f"{name}_{h}.{fmt}"
    dst = figures_dir / fname
    if not dst.exists():
        dst.write_bytes(data)
    return str(dst)


def extract_pdf_figures(doc, pages_dir: Path, out: Path) -> list[dict]:
    """Extract embedded images from each PDF page."""
    import fitz
    figures_dir = out / "figures"
    figures_dir.mkdir(parents=True, exist_ok=True)
    entries = []
    for i in range(len(doc)):
        page = doc[i]
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base = doc.extract_image(xref)
            ext = base["ext"]
            data = base["image"]
            path = _save_figure(data, f"p{i+1:03d}_img{img_index}",
                                figures_dir, ext)
            entries.append({
                "page": i + 1,
                "index": img_index,
                "file": path,
                "format": ext,
                "width": base.get("width"),
                "height": base.get("height"),
                "source": "pdf_embedded"
            })
    return entries


def extract_pdf(src: Path, out: Path, dpi: int, extract_figures: bool = False) -> dict:
    import fitz  # PyMuPDF

    doc = fitz.open(src)
    pages_dir = out / "pages"
    pages_dir.mkdir(parents=True, exist_ok=True)
    manifest = {"source": str(src), "type": "pdf", "page_count": len(doc), "pages": []}
    text_parts = []
    zoom = dpi / 72.0
    figure_entries = []

    if extract_figures:
        figure_entries = extract_pdf_figures(doc, pages_dir, out)

    for i, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        img_count = len(page.get_images(full=True))
        draw_count = len(page.get_drawings())
        png = pages_dir / f"p{i:03d}.png"
        page.get_pixmap(matrix=fitz.Matrix(zoom, zoom)).save(png)
        flags = []
        if len(text) < 200:
            flags.append("LOW_TEXT")  # likely scanned or diagram-heavy: visual read REQUIRED
        if img_count > 0 or draw_count > 8:
            flags.append("HAS_GRAPHICS")  # contains figures: visual read REQUIRED
        manifest["pages"].append(
            {"page": i, "chars": len(text), "images": img_count,
             "vector_drawings": draw_count, "png": str(png), "flags": flags}
        )
        text_parts.append(f"\n\n===== PAGE {i} / {len(doc)} =====\n{text}")
    (out / "text.md").write_text("".join(text_parts), encoding="utf-8")

    if figure_entries:
        manifest["figures"] = figure_entries
        (out / "figures.json").write_text(
            json.dumps(figure_entries, indent=2), encoding="utf-8")
        print(f"  Extracted {len(figure_entries)} embedded figures to figures/")

    doc.close()
    return manifest


def export_pptx_images(src: Path, pages_dir: Path, count_hint: int) -> bool:
    """Render real slide images via PowerPoint COM (full fidelity)."""
    try:
        import win32com.client
        app = win32com.client.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open(str(src.resolve()), ReadOnly=True,
                                      Untitled=False, WithWindow=False)
        try:
            for i in range(1, pres.Slides.Count + 1):
                png = pages_dir / f"p{i:03d}.png"
                pres.Slides(i).Export(str(png.resolve()), "PNG", 1600, 900)
        finally:
            pres.Close()
            app.Quit()
        return True
    except Exception as e:  # noqa: BLE001 - report and fall back to text-only
        print(f"WARNING: COM slide export failed ({e}); images unavailable", file=sys.stderr)
        return False


def extract_pptx_figures(pres, out: Path) -> list[dict]:
    """Extract embedded images from PPTX slides."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    figures_dir = out / "figures"
    figures_dir.mkdir(parents=True, exist_ok=True)
    entries = []
    for i, slide in enumerate(pres.slides, start=1):
        for idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                ext = img.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                data = img.blob
                path = _save_figure(data, f"p{i:03d}_fig{idx}",
                                    figures_dir, ext)
                entries.append({
                    "page": i,
                    "index": idx,
                    "file": path,
                    "format": ext,
                    "width": shape.width,
                    "height": shape.height,
                    "source": "pptx_picture"
                })
    return entries


def extract_pptx(src: Path, out: Path, extract_figures: bool = False) -> dict:
    from pptx import Presentation

    pres = Presentation(src)
    pages_dir = out / "pages"
    pages_dir.mkdir(parents=True, exist_ok=True)
    slides = list(pres.slides)
    manifest = {"source": str(src), "type": "pptx", "page_count": len(slides), "pages": []}
    text_parts = []
    figure_entries = []

    if extract_figures:
        figure_entries = extract_pptx_figures(pres, out)

    for i, slide in enumerate(slides, start=1):
        chunks = []
        shape_kinds = []
        for shape in slide.shapes:
            shape_kinds.append(str(shape.shape_type))
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    chunks.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.append(" | ".join(c.text.strip() for c in row.cells))
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        body = "\n".join(chunks)
        flags = []
        if len(body) < 100:
            flags.append("LOW_TEXT")
        if any("PICTURE" in k or "FREEFORM" in k or "GROUP" in k or "AUTO_SHAPE" in k
               for k in shape_kinds):
            flags.append("HAS_GRAPHICS")
        manifest["pages"].append(
            {"page": i, "chars": len(body), "shapes": shape_kinds,
             "png": str(pages_dir / f"p{i:03d}.png"), "flags": flags}
        )
        part = f"\n\n===== SLIDE {i} / {len(slides)} =====\n{body}"
        if notes:
            part += f"\n--- SPEAKER NOTES ---\n{notes}"
        text_parts.append(part)
    (out / "text.md").write_text("".join(text_parts), encoding="utf-8")

    if figure_entries:
        manifest["figures"] = figure_entries
        (out / "figures.json").write_text(
            json.dumps(figure_entries, indent=2), encoding="utf-8")
        print(f"  Extracted {len(figure_entries)} embedded figures to figures/")

    manifest["slide_images_exported"] = export_pptx_images(src, pages_dir, len(slides))
    return manifest


def extract_docx_figures(src: Path, out: Path) -> list[dict]:
    """Extract embedded images from DOCX."""
    import zipfile
    from lxml import etree

    figures_dir = out / "figures"
    figures_dir.mkdir(parents=True, exist_ok=True)
    entries = []

    nsmap = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    with zipfile.ZipFile(src) as z:
        # Find image relationships
        rels = {}
        for name in z.namelist():
            if name.startswith("word/_rels/") and name.endswith(".rels"):
                rel_tree = etree.fromstring(z.read(name))
                for rel_elem in rel_tree:
                    rid = rel_elem.get("Id")
                    target = rel_elem.get("Target")
                    if target and ("media" in target.lower() or "image" in target.lower()):
                        rels[rid] = target

        # Walk document.xml for image references
        if "word/document.xml" in z.namelist():
            tree = etree.fromstring(z.read("word/document.xml"))
            blip_fills = tree.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}blip")
            page = 1  # DOCX has no fixed page breaks
            for idx, blip in enumerate(blip_fills):
                embed = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                if embed and embed in rels:
                    img_path = "word/" + rels[embed]
                    if img_path in z.namelist():
                        data = z.read(img_path)
                        ext = img_path.rsplit(".", 1)[-1]
                        if ext.lower() in ("png", "jpg", "jpeg", "gif", "bmp", "tiff"):
                            path = _save_figure(data, f"docx_fig{idx}",
                                                figures_dir, ext)
                            entries.append({
                                "page": page,
                                "index": idx,
                                "file": path,
                                "format": ext,
                                "source": "docx_embedded"
                            })
    return entries


def extract_docx(src: Path, out: Path, extract_figures: bool = False) -> dict:
    import docx

    d = docx.Document(src)
    parts = [p.text for p in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    text = "\n".join(parts)
    (out / "text.md").write_text(text, encoding="utf-8")

    manifest = {"source": str(src), "type": "docx", "page_count": None,
                "pages": [], "chars": len(text)}

    figure_entries = []
    if extract_figures:
        figure_entries = extract_docx_figures(src, out)
        if figure_entries:
            manifest["figures"] = figure_entries
            (out / "figures.json").write_text(
                json.dumps(figure_entries, indent=2), encoding="utf-8")
            print(f"  Extracted {len(figure_entries)} embedded figures to figures/")

    return manifest


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("input")
    ap.add_argument("--outdir", default="_extracted")
    ap.add_argument("--dpi", type=int, default=150)
    ap.add_argument("--extract-figures", action="store_true",
                    help="Extract embedded images from PDF/PPTX/DOCX into figures/")
    args = ap.parse_args()

    src = Path(args.input)
    if not src.exists():
        print(f"ERROR: not found: {src}", file=sys.stderr)
        return 1
    out = Path(args.outdir) / src.stem
    out.mkdir(parents=True, exist_ok=True)

    ext = src.suffix.lower()
    if ext == ".pdf":
        manifest = extract_pdf(src, out, args.dpi, args.extract_figures)
    elif ext in (".pptx", ".ppt"):
        manifest = extract_pptx(src, out, args.extract_figures)
    elif ext == ".docx":
        manifest = extract_docx(src, out, args.extract_figures)
    else:
        print(f"ERROR: unsupported type: {ext}", file=sys.stderr)
        return 1

    must_view = [p["page"] for p in manifest["pages"] if p["flags"]]
    manifest["pages_requiring_visual_read"] = must_view
    mpath = out / "manifest.json"
    mpath.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(f"Extracted {manifest.get('page_count')} pages/slides.")
    print(f"Pages REQUIRING visual read (graphics or low text): {must_view}")
    if args.extract_figures:
        nf = len(manifest.get("figures", []))
        print(f"Figures extracted: {nf} (see figures/ and figures.json)")
    print(str(mpath))
    return 0


if __name__ == "__main__":
    sys.exit(main())
